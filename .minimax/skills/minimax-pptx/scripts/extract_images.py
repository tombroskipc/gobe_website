#!/usr/bin/env python3
"""Extract images from a PPTX file with slide-level attribution.

Usage:
    python extract_images.py presentation.pptx
    python extract_images.py presentation.pptx --output-dir ./images
    python extract_images.py presentation.pptx --slides 1,3,5
    python extract_images.py presentation.pptx --range 2-4
    python extract_images.py presentation.pptx --list-only

Requires: python-pptx  (pip install python-pptx)
Works with Python >= 3.7.
"""

from __future__ import annotations

import argparse
import os
import re
import sys
from dataclasses import dataclass

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class ImageInfo:
    slide_number: int
    shape_name: str
    content_type: str
    extension: str
    width: int | None
    height: int | None
    blob: bytes


def extract_images_from_slide(slide, slide_number: int) -> list[ImageInfo]:
    """Extract all images from a single slide, including grouped shapes."""
    images: list[ImageInfo] = []

    def process_shape(shape):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                process_shape(child)
            return

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img = shape.image
            images.append(
                ImageInfo(
                    slide_number=slide_number,
                    shape_name=shape.name,
                    content_type=img.content_type,
                    extension=img.ext,
                    width=shape.width,
                    height=shape.height,
                    blob=img.blob,
                )
            )

    for shape in slide.shapes:
        process_shape(shape)

    return images


def load_images(pptx_path: str) -> list[ImageInfo]:
    presentation = Presentation(pptx_path)
    all_images: list[ImageInfo] = []
    for index, slide in enumerate(presentation.slides, 1):
        all_images.extend(extract_images_from_slide(slide, index))
    return all_images


def filter_by_slides(
    images: list[ImageInfo], slide_numbers: set[int]
) -> list[ImageInfo]:
    return [img for img in images if img.slide_number in slide_numbers]


def filter_by_range(images: list[ImageInfo], start: int, end: int) -> list[ImageInfo]:
    return [img for img in images if start <= img.slide_number <= end]


def parse_slide_numbers(raw_value: str) -> set[int]:
    numbers: set[int] = set()
    for chunk in raw_value.split(","):
        value = chunk.strip()
        if not value:
            continue
        if not value.isdigit():
            raise ValueError(f'Invalid slide number "{value}".')
        slide_number = int(value)
        if slide_number < 1:
            raise ValueError("Slide numbers must be >= 1.")
        numbers.add(slide_number)
    if not numbers:
        raise ValueError("At least one slide number must be provided.")
    return numbers


def parse_slide_range(raw_value: str) -> tuple[int, int]:
    match = re.fullmatch(r"\s*(\d+)\s*-\s*(\d+)\s*", raw_value)
    if not match:
        raise ValueError('Invalid range format. Use "start-end", e.g. "2-4".')
    start, end = int(match.group(1)), int(match.group(2))
    if start < 1 or end < 1:
        raise ValueError("Range values must be >= 1.")
    if start > end:
        raise ValueError("Range start must be <= end.")
    return start, end


def emu_to_inches(emu: int | None) -> str:
    if emu is None:
        return "?"
    return f"{emu / 914400:.1f}"


def print_image_list(images: list[ImageInfo]) -> None:
    if not images:
        print("No images found.")
        return

    for img in images:
        width = emu_to_inches(img.width)
        height = emu_to_inches(img.height)
        size_kb = len(img.blob) / 1024
        print(
            f"Slide {img.slide_number:>3} | {img.shape_name:<30} | "
            f"{img.content_type:<20} | {width}x{height} in | {size_kb:.0f} KB"
        )

    print(f"\nTotal: {len(images)} image(s)")


def save_images(images: list[ImageInfo], output_dir: str) -> None:
    if not images:
        print("No images found.")
        return

    os.makedirs(output_dir, exist_ok=True)
    slide_counters: dict[int, int] = {}

    for img in images:
        count = slide_counters.get(img.slide_number, 0) + 1
        slide_counters[img.slide_number] = count

        safe_name = re.sub(r"[^\w\-.]", "_", img.shape_name)
        filename = (
            f"slide{img.slide_number:02d}_{count:02d}_{safe_name}.{img.extension}"
        )
        filepath = os.path.join(output_dir, filename)

        with open(filepath, "wb") as file_obj:
            file_obj.write(img.blob)

        width = emu_to_inches(img.width)
        height = emu_to_inches(img.height)
        size_kb = len(img.blob) / 1024
        print(f"Saved: {filename} ({width}x{height} in, {size_kb:.0f} KB)")

    print(f"\n{len(images)} image(s) saved to {output_dir}/")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Extract images from a PPTX file with slide attribution."
    )
    parser.add_argument("pptx_path", help="Path to the PPTX file.")
    parser.add_argument(
        "--output-dir",
        "-o",
        default="./extracted_images",
        help="Directory to save extracted images (default: ./extracted_images).",
    )
    selection = parser.add_mutually_exclusive_group()
    selection.add_argument(
        "--slides",
        help='Only extract from these slides, e.g. "1,3,5".',
    )
    selection.add_argument(
        "--range",
        dest="slide_range",
        help='Only extract from this slide range, e.g. "2-4".',
    )
    parser.add_argument(
        "--list-only",
        "-l",
        action="store_true",
        help="List images with metadata instead of saving them.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()

    try:
        images = load_images(args.pptx_path)

        if args.slides:
            numbers = parse_slide_numbers(args.slides)
            images = filter_by_slides(images, numbers)
        elif args.slide_range:
            start, end = parse_slide_range(args.slide_range)
            images = filter_by_range(images, start, end)

        if args.list_only:
            print_image_list(images)
        else:
            save_images(images, args.output_dir)

    except ValueError as error:
        print(str(error), file=sys.stderr)
        return 1
    except Exception as error:
        print(f"Failed to extract images: {error}", file=sys.stderr)
        return 1

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
