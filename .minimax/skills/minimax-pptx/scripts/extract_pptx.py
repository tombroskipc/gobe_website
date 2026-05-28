#!/usr/bin/env python3
"""Extract or search PPTX text content, slide by slide.

Usage:
    python extract_pptx.py <path_to_pptx>
    python extract_pptx.py <path_to_pptx> --slides 1,3,5
    python extract_pptx.py <path_to_pptx> --range 2-4
    python extract_pptx.py <path_to_pptx> --search QUERY [--ignore-case]
    python extract_pptx.py <path_to_pptx> --search PATTERN --regex [--ignore-case]
    python extract_pptx.py <path_to_pptx> --notes

Requires: python-pptx  (pip install python-pptx)
Works with Python >= 3.7 (no Python 3.10+ requirement).
"""

from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass

from pptx import Presentation


@dataclass
class SlideContent:
    slide_number: int
    lines: list[str]
    notes: str | None = None


def collect_slide_notes(slide) -> str | None:
    """Extract speaker notes text from a slide, if present."""
    if not slide.has_notes_slide:
        return None
    notes_text = slide.notes_slide.notes_text_frame.text.strip()
    return notes_text if notes_text else None


def collect_slide_lines(slide) -> list[str]:
    lines: list[str] = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    lines.append(text)

        if shape.has_table:
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    lines.append(" | ".join(cells))

    return lines


def load_presentation_content(pptx_path: str) -> list[SlideContent]:
    presentation = Presentation(pptx_path)
    return [
        SlideContent(
            slide_number=index,
            lines=collect_slide_lines(slide),
            notes=collect_slide_notes(slide),
        )
        for index, slide in enumerate(presentation.slides, 1)
    ]


def build_matcher(query: str, use_regex: bool, ignore_case: bool):
    if use_regex:
        flags = re.IGNORECASE if ignore_case else 0
        pattern = re.compile(query, flags)
        return pattern.search

    normalized_query = query.casefold() if ignore_case else query

    def contains(text: str):
        haystack = text.casefold() if ignore_case else text
        return normalized_query in haystack

    return contains


def parse_slide_numbers(raw_value: str) -> set[int]:
    slide_numbers: set[int] = set()

    for chunk in raw_value.split(","):
        value = chunk.strip()
        if not value:
            continue
        if not value.isdigit():
            raise ValueError(
                f'Invalid slide number "{value}". Use comma-separated positive integers.'
            )

        slide_number = int(value)
        if slide_number < 1:
            raise ValueError("Slide numbers must be >= 1.")
        slide_numbers.add(slide_number)

    if not slide_numbers:
        raise ValueError("At least one slide number must be provided.")

    return slide_numbers


def parse_slide_range(raw_value: str) -> tuple[int, int]:
    match = re.fullmatch(r"\s*(\d+)\s*-\s*(\d+)\s*", raw_value)
    if not match:
        raise ValueError('Invalid range format. Use "start-end", for example "2-4".')

    start = int(match.group(1))
    end = int(match.group(2))
    if start < 1 or end < 1:
        raise ValueError("Slide range values must be >= 1.")
    if start > end:
        raise ValueError("Slide range start must be <= end.")

    return start, end


def select_slides(
    slides: list[SlideContent],
    selected_numbers: set[int] | None,
    selected_range: tuple[int, int] | None,
) -> list[SlideContent]:
    total_slides = len(slides)

    if selected_numbers is not None:
        invalid = sorted(
            slide_number
            for slide_number in selected_numbers
            if slide_number > total_slides
        )
        if invalid:
            raise ValueError(
                f"Requested slide(s) out of range: {', '.join(str(value) for value in invalid)}. "
                f"Presentation has {total_slides} slide(s)."
            )
        return [slide for slide in slides if slide.slide_number in selected_numbers]

    if selected_range is not None:
        start, end = selected_range
        if start > total_slides:
            raise ValueError(
                f"Requested range {start}-{end} is out of range. Presentation has {total_slides} slide(s)."
            )
        if end > total_slides:
            raise ValueError(
                f"Requested range end {end} is out of range. Presentation has {total_slides} slide(s)."
            )
        return [slide for slide in slides if start <= slide.slide_number <= end]

    return slides


def print_slide_content(slides: list[SlideContent], show_notes: bool) -> None:
    for slide in slides:
        print(f"=== Slide {slide.slide_number} ===")
        for line in slide.lines:
            print(line)
        if show_notes and slide.notes:
            print("[Speaker Notes]")
            print(slide.notes)
        print()


def print_search_results(
    slides: list[SlideContent],
    query: str,
    use_regex: bool,
    ignore_case: bool,
    show_notes: bool,
) -> None:
    matcher = build_matcher(query, use_regex, ignore_case)
    matched_slides = 0
    matched_lines = 0

    for slide in slides:
        matches = [line for line in slide.lines if matcher(line)]
        notes_match = show_notes and slide.notes and matcher(slide.notes)
        if not matches and not notes_match:
            continue

        matched_slides += 1
        matched_lines += len(matches)
        if notes_match:
            matched_lines += 1
        print(f"=== Slide {slide.slide_number} ===")
        for line in matches:
            print(line)
        if notes_match:
            print("[Speaker Notes]")
            print(slide.notes)
        print()

    if matched_lines == 0:
        print(f'No matches found for "{query}".')
        return

    print(
        f'Found {matched_lines} matching line(s) across {matched_slides} slide(s).'
    )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Extract all PPTX text or search matching content."
    )
    parser.add_argument("pptx_path", help="Path to the PPTX file to inspect.")
    selection_group = parser.add_mutually_exclusive_group()
    selection_group.add_argument(
        "--slides",
        help='Only inspect these slide numbers, for example "1,3,5".',
    )
    selection_group.add_argument(
        "--range",
        dest="slide_range",
        help='Only inspect a continuous slide range, for example "2-4".',
    )
    parser.add_argument(
        "--search",
        help="Only print lines that match this query instead of the full deck.",
    )
    parser.add_argument(
        "--regex",
        action="store_true",
        help="Treat --search as a regular expression.",
    )
    parser.add_argument(
        "--ignore-case",
        action="store_true",
        help="Make search case-insensitive.",
    )
    parser.add_argument(
        "--notes",
        action="store_true",
        help="Include speaker notes in the output.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()

    if args.regex and not args.search:
        print("--regex requires --search.", file=sys.stderr)
        return 1

    try:
        slides = load_presentation_content(args.pptx_path)
        selected_numbers = parse_slide_numbers(args.slides) if args.slides else None
        selected_range = (
            parse_slide_range(args.slide_range) if args.slide_range else None
        )
        slides = select_slides(slides, selected_numbers, selected_range)

        if args.search:
            print_search_results(
                slides,
                query=args.search,
                use_regex=args.regex,
                ignore_case=args.ignore_case,
                show_notes=args.notes,
            )
        else:
            print_slide_content(slides, show_notes=args.notes)
    except ValueError as error:
        print(str(error), file=sys.stderr)
        return 1
    except re.error as error:
        print(f"Invalid regular expression: {error}", file=sys.stderr)
        return 1
    except Exception as error:
        print(f"Failed to inspect PPTX: {error}", file=sys.stderr)
        return 1

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
